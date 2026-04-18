#!/usr/bin/env python3
"""
Generate NeoTherm / Global-Therma pitch deck and appendix as PPTX files.
Outputs:
    deck.pptx     — 19 main pitch slides
    appendix.pptx — 9 appendix slides
"""

import os
import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ---------------------------------------------------------------------------
# Directory
# ---------------------------------------------------------------------------
BASE = os.path.dirname(os.path.abspath(__file__))

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
C = {
    "ink":        RGBColor(0x1a, 0x23, 0x32),
    "ink_light":  RGBColor(0x3d, 0x4f, 0x63),
    "ink_muted":  RGBColor(0x64, 0x74, 0x8b),
    "surface":    RGBColor(0xf8, 0xf9, 0xfb),
    "white":      RGBColor(0xff, 0xff, 0xff),
    "border":     RGBColor(0xdd, 0xe2, 0xe8),
    "teal_deep":  RGBColor(0x0c, 0x6e, 0x6e),
    "teal":       RGBColor(0x0d, 0x8a, 0x8a),
    "teal_light": RGBColor(0xe6, 0xf5, 0xf5),
    "blue":       RGBColor(0x25, 0x63, 0xa0),
    "amber":      RGBColor(0xb8, 0x86, 0x0b),
    "red":        RGBColor(0xc0, 0x39, 0x2b),
    "red_light":  RGBColor(0xfd, 0xf0, 0xee),
    "green":      RGBColor(0x1a, 0x7a, 0x4c),
    "green_light":RGBColor(0xed, 0xf7, 0xf1),
    "phase1":     RGBColor(0x25, 0x63, 0xa0),
    "phase2":     RGBColor(0x0d, 0x8a, 0x8a),
    "phase3":     RGBColor(0x6b, 0x4c, 0x9a),
    "hack_blue":  RGBColor(0x15, 0x65, 0xC0),
    "orange":     RGBColor(0xd3, 0x54, 0x00),
    "light_red":  RGBColor(0xe7, 0x4c, 0x3c),
}

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _set_font(run, *, size=None, color=None, bold=False, italic=False,
              name="Calibri", underline=False):
    f = run.font
    if size:
        f.size = Pt(size)
    if color:
        f.color.rgb = color
    f.bold = bold
    f.italic = italic
    f.name = name
    f.underline = underline


def _add_textbox(slide, left, top, width, height, text, *,
                 size=14, color=None, bold=False, italic=False,
                 name="Calibri", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 word_wrap=True):
    """Convenience: add a textbox with a single run."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tb.text_frame.word_wrap = word_wrap
    tf = tb.text_frame
    tf.paragraphs[0].alignment = align
    try:
        tf.paragraphs[0].vertical_anchor = anchor
    except Exception:
        pass
    r = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run()
    r.text = text
    _set_font(r, size=size, color=color or C["ink"], bold=bold, italic=italic, name=name)
    return tb


def _add_rich_textbox(slide, left, top, width, height, parts, *,
                      align=PP_ALIGN.LEFT, line_spacing=None, space_after=None):
    """parts = list of (text, dict-of-font-kwargs)."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    if space_after:
        p.space_after = Pt(space_after)
    for txt, kw in parts:
        r = p.add_run()
        r.text = txt
        _set_font(r, **kw)
    return tb


def _add_multiline_textbox(slide, left, top, width, height, lines, *,
                           default_align=PP_ALIGN.LEFT):
    """lines = list of (text, dict-of-font-kwargs, optional-alignment)."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tb.text_frame.word_wrap = True
    for i, item in enumerate(lines):
        text = item[0]
        kw = item[1] if len(item) > 1 else {}
        al = item[2] if len(item) > 2 else default_align
        if i == 0:
            p = tb.text_frame.paragraphs[0]
        else:
            p = tb.text_frame.add_paragraph()
        p.alignment = al
        r = p.add_run()
        r.text = text
        _set_font(r, **kw)
        if "space_after" in kw:
            p.space_after = Pt(kw.pop("space_after"))
    return tb


def _rect(slide, left, top, w, h, *, fill=None, border_color=None, border_w=None,
          shape=MSO_SHAPE.RECTANGLE, rotation=0):
    s = slide.shapes.add_shape(shape, Inches(left), Inches(top),
                                Inches(w), Inches(h))
    s.rotation = rotation
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = Pt(border_w or 1)
    else:
        s.line.fill.background()
    return s


def _rounded_rect(slide, left, top, w, h, *, fill=None, border_color=None, border_w=None):
    return _rect(slide, left, top, w, h, fill=fill, border_color=border_color,
                 border_w=border_w, shape=MSO_SHAPE.ROUNDED_RECTANGLE)


def _card(slide, left, top, w, h, *, fill=None, border_color=None,
          top_border_color=None, top_border_h=0.06):
    """Card = rounded rect + optional coloured top bar."""
    card = _rounded_rect(slide, left, top, w, h,
                         fill=fill or C["surface"],
                         border_color=border_color or C["border"])
    if top_border_color:
        _rect(slide, left + 0.02, top + 0.02, w - 0.04, top_border_h,
              fill=top_border_color)
    return card


def _arrow_right(slide, left, top, w=0.5, h=0.3, fill=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left), Inches(top),
                                Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill or C["teal"]
    s.line.fill.background()
    return s


def _accent_bar(slide, color=None):
    """Thin teal bar across top of a white slide."""
    _rect(slide, 0, 0, 13.333, 0.06, fill=color or C["teal"])


def _label(slide, left, top, text, color=None):
    """Section label (uppercase, small, spaced)."""
    tb = _add_textbox(slide, left, top, 6, 0.3, text,
                      size=10, color=color or C["teal"], bold=True, name="Calibri")
    return tb


def _title(slide, left, top, text, *, width=11, size=36, color=None, name="Calibri Light"):
    return _add_textbox(slide, left, top, width, 1.0, text,
                        size=size, color=color or C["ink"], bold=False, name=name)


def _subtitle(slide, left, top, text, *, width=11, size=18, color=None):
    return _add_textbox(slide, left, top, width, 0.6, text,
                        size=size, color=color or C["ink_light"], name="Calibri")


def _solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _img(slide, path, left, top, width=None, height=None):
    kw = {}
    if width: kw["width"] = Inches(width)
    if height: kw["height"] = Inches(height)
    return slide.shapes.add_picture(os.path.join(BASE, path),
                                    Inches(left), Inches(top), **kw)


def _new_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def _code_block(slide, left, top, w, h, lines, *, bg=None):
    """Dark code block with coloured monospace text.
    lines = list of (text, color) tuples."""
    _rounded_rect(slide, left, top, w, h, fill=bg or C["ink"])
    y = top + 0.2
    for txt, clr in lines:
        _add_textbox(slide, left + 0.3, y, w - 0.6, 0.3, txt,
                     size=11, color=clr, name="Consolas")
        y += 0.26
    return y


def _stat_box(slide, left, top, w, h, number, label, *, num_color=None,
              bg=None, num_size=28, label_size=11):
    _rounded_rect(slide, left, top, w, h, fill=bg or C["surface"],
                  border_color=C["border"])
    _add_textbox(slide, left + 0.1, top + 0.1, w - 0.2, 0.5, number,
                 size=num_size, color=num_color or C["teal"], bold=True,
                 align=PP_ALIGN.CENTER)
    _add_textbox(slide, left + 0.1, top + 0.55, w - 0.2, 0.4, label,
                 size=label_size, color=C["ink_muted"], align=PP_ALIGN.CENTER)


def _small_stat_box(slide, left, top, w, h, number, label, *, num_color=None,
                    bg=None, num_size=22, label_size=10):
    _rounded_rect(slide, left, top, w, h, fill=bg or C["surface"],
                  border_color=C["border"])
    _add_textbox(slide, left + 0.05, top + 0.05, w - 0.1, 0.35, number,
                 size=num_size, color=num_color or C["teal"], bold=True,
                 align=PP_ALIGN.CENTER)
    _add_textbox(slide, left + 0.05, top + 0.38, w - 0.1, 0.35, label,
                 size=label_size, color=C["ink_muted"], align=PP_ALIGN.CENTER)


# ===========================================================================
# MAIN DECK  (19 slides)
# ===========================================================================

def build_deck():
    prs = _new_prs()

    # ---- Slide 1: Blue Opener ------------------------------------------------
    s = _blank(prs)
    _solid_bg(s, C["hack_blue"])
    _add_textbox(s, 1, 2.2, 11.333, 1.0, "Team 218 \u2014 Global-Therma",
                 size=36, color=C["white"], bold=True, align=PP_ALIGN.CENTER)
    _add_textbox(s, 1.5, 3.5, 10.333, 1.2,
                 "Infrared thermography-enabled NICU monitoring\nfor early detection of circulatory deterioration",
                 size=18, color=RGBColor(0xB3, 0xCC, 0xFF), align=PP_ALIGN.CENTER)
    _add_textbox(s, 2, 5.8, 9.333, 0.5, "\u2705  InterSystems Prize",
                 size=16, color=C["white"], align=PP_ALIGN.CENTER, bold=True)

    # ---- Slide 2: Baby Photo -------------------------------------------------
    s = _blank(prs)
    try:
        pic = _img(s, "india_nicu-dfid-crop.jpg", 0, 0, width=13.333, height=7.5)
    except Exception:
        _rect(s, 0, 0, 13.333, 7.5, fill=C["ink"])
    # dark gradient overlay at bottom
    _rect(s, 0, 5.0, 13.333, 2.5, fill=RGBColor(0x00, 0x00, 0x00))
    # make it semi-transparent by setting alpha (pptx doesn't directly support, but we overlay)
    overlay = s.shapes[-1]
    from pptx.oxml.ns import qn
    solidFill = overlay.fill._fill
    # set transparency
    try:
        srgb = solidFill.find(qn("a:solidFill")).find(qn("a:srgbClr"))
        if srgb is None:
            srgb = solidFill.find(qn("a:srgbClr"))
        if srgb is not None:
            from lxml import etree
            alpha = etree.SubElement(srgb, qn("a:alpha"))
            alpha.set("val", "60000")  # 60% opacity
    except Exception:
        pass
    _add_textbox(s, 1.0, 5.5, 11.333, 1.5,
                 "Look at him. Just days old, and already fighting for his life.",
                 size=36, color=C["white"], italic=True, name="Georgia",
                 align=PP_ALIGN.CENTER)

    # ---- Slide 3: Globe Pie Chart --------------------------------------------
    s = _blank(prs)
    _accent_bar(s)
    _label(s, 0.8, 0.3, "GLOBAL NEONATAL DEATHS FROM CIRCULATORY FAILURE")

    # Doughnut chart via python-pptx
    chart_data = CategoryChartData()
    chart_data.categories = ["Late-onset sepsis", "NEC", "Early-onset sepsis",
                             "Cardiogenic shock", "PDA / other"]
    chart_data.add_series("Deaths", (34, 22, 18, 12, 14))
    chart_frame = s.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, Inches(0.8), Inches(1.2),
        Inches(5.5), Inches(5.0), chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    # colour each slice
    series = chart.series[0]
    slice_colors = [C["red"], C["orange"], C["light_red"], C["blue"], C["ink_muted"]]
    for i, clr in enumerate(slice_colors):
        pt = series.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = clr

    # Center hero number
    _add_textbox(s, 1.5, 2.8, 4.0, 0.8, "680,000",
                 size=40, color=C["ink"], bold=True, align=PP_ALIGN.CENTER)
    _add_textbox(s, 1.5, 3.6, 4.0, 0.5, "deaths per year",
                 size=16, color=C["ink_muted"], align=PP_ALIGN.CENTER)

    # Legend on the right
    legend_items = [
        ("Late-onset sepsis", "34%", C["red"]),
        ("NEC", "22%", C["orange"]),
        ("Early-onset sepsis", "18%", C["light_red"]),
        ("Cardiogenic shock", "12%", C["blue"]),
        ("PDA / other", "14%", C["ink_muted"]),
    ]
    ly = 1.8
    for label_text, pct, clr in legend_items:
        _rect(s, 7.2, ly + 0.05, 0.25, 0.25, fill=clr)
        _add_textbox(s, 7.6, ly, 3.5, 0.35, f"{label_text}  {pct}",
                     size=14, color=C["ink"])
        ly += 0.55

    # Bottom callout
    _add_textbox(s, 0.8, 6.5, 11.5, 0.5,
                 "NeoTherm detects the thermal signature common to all five.",
                 size=16, color=C["teal"], bold=True, align=PP_ALIGN.CENTER)

    # ---- Slide 4: Problem: Detection Fails -----------------------------------
    s = _blank(prs)
    _accent_bar(s)
    _label(s, 0.8, 0.3, "THE PROBLEM")
    _title(s, 0.8, 0.7, "Current diagnostics are nonspecific, variable, or late.")

    cards = [
        (C["red"], "\u23F1  Too Late",
         "Blood cultures take 24\u201372 hours. Clinical signs appear only after deterioration. Spot checks every 1\u20134 hours miss what happens between readings."),
        (C["amber"], "\u26A0  Too Invasive",
         "Adhesive probes damage premature skin. Blood draws for CRP / procalcitonin. 78% of NICU devices used off-label."),
        (C["ink_muted"], "\u25C9  Too Nonspecific",
         "Single-point probes can\u2019t detect spatial gradients. Capillary refill unreliable on dark skin. qSOFA designed for adults."),
    ]
    cx = 0.8
    for border_clr, title_text, body_text in cards:
        _card(s, cx, 2.0, 3.7, 3.2, top_border_color=border_clr, top_border_h=0.08)
        _add_textbox(s, cx + 0.25, 2.25, 3.2, 0.5, title_text,
                     size=22, color=C["ink"], bold=True)
        _add_textbox(s, cx + 0.25, 2.85, 3.2, 2.2, body_text,
                     size=14, color=C["ink_light"])
        cx += 3.95

    # Red bottom bar
    _rect(s, 0.8, 5.6, 11.733, 0.7, fill=C["red_light"])
    _add_textbox(s, 1.0, 5.65, 7, 0.6,
                 "The gap: No continuous, non-contact, automated screening.",
                 size=15, color=C["red"], bold=True)
    _add_textbox(s, 9.5, 5.65, 3.0, 0.6, "$2.1B / yr",
                 size=26, color=C["red"], bold=True, align=PP_ALIGN.RIGHT)

    # ---- Slide 5: "But what if..." -------------------------------------------
    s = _blank(prs)
    _accent_bar(s)
    _add_textbox(s, 1.5, 1.8, 10.333, 0.6,
                 "What if we could see the warning signs 5 hours earlier?",
                 size=16, color=C["ink_muted"], align=PP_ALIGN.CENTER)
    _add_textbox(s, 1.5, 2.8, 10.333, 1.2, "But what if\u2026",
                 size=60, color=C["ink"], name="Georgia", align=PP_ALIGN.CENTER)
    # Dashed placeholder
    placeholder = _rounded_rect(s, 2.5, 4.2, 8.333, 2.6,
                                border_color=C["border"], border_w=2)
    placeholder.fill.background()
    _add_textbox(s, 3.0, 4.8, 7.333, 0.6,
                 "[ Claude GIF / Demo Animation here ]",
                 size=18, color=C["ink_muted"], align=PP_ALIGN.CENTER)
    _add_textbox(s, 3.0, 5.5, 7.333, 0.4,
                 "Replace in Google Slides",
                 size=12, color=C["border"], align=PP_ALIGN.CENTER)

    # ---- Slide 6: CPTD Science -----------------------------------------------
    s = _blank(prs)
    _accent_bar(s)
    _label(s, 0.8, 0.3, "THE SCIENCE")
    _title(s, 0.8, 0.7, "One camera. The body tells us everything.")

    # Left side: simplified body diagram
    # Head
    _rect(s, 2.5, 1.8, 1.2, 1.0, fill=RGBColor(0xFF, 0xCC, 0xCC),
          shape=MSO_SHAPE.OVAL)
    # Torso
    _rounded_rect(s, 2.2, 2.9, 1.8, 1.8, fill=RGBColor(0xFF, 0x99, 0x99))
    _add_textbox(s, 2.3, 3.4, 1.6, 0.5, "CORE\n36.8\u00b0C",
                 size=11, color=C["white"], bold=True, align=PP_ALIGN.CENTER, name="Consolas")
    # Left arm
    _rect(s, 1.2, 3.0, 0.8, 0.4, fill=RGBColor(0xAD, 0xD8, 0xE6))
    # Left hand
    _rect(s, 0.7, 3.0, 0.5, 0.45, fill=RGBColor(0x73, 0xB0, 0xD7),
          shape=MSO_SHAPE.OVAL)
    _add_textbox(s, 0.5, 3.5, 1.0, 0.3, "33.1\u00b0C",
                 size=9, color=C["blue"], bold=True, align=PP_ALIGN.CENTER, name="Consolas")
    # Right arm
    _rect(s, 4.2, 3.0, 0.8, 0.4, fill=RGBColor(0xAD, 0xD8, 0xE6))
    # Right hand
    _rect(s, 4.9, 3.0, 0.5, 0.45, fill=RGBColor(0x73, 0xB0, 0xD7),
          shape=MSO_SHAPE.OVAL)
    _add_textbox(s, 4.7, 3.5, 1.0, 0.3, "33.6\u00b0C",
                 size=9, color=C["blue"], bold=True, align=PP_ALIGN.CENTER, name="Consolas")
    # Left leg
    _rect(s, 2.3, 4.8, 0.5, 1.2, fill=RGBColor(0xAD, 0xD8, 0xE6))
    # Left foot
    _rect(s, 2.2, 6.0, 0.6, 0.4, fill=RGBColor(0x4A, 0x90, 0xC4),
          shape=MSO_SHAPE.OVAL)
    _add_textbox(s, 1.8, 6.35, 1.2, 0.3, "31.5\u00b0C",
                 size=9, color=C["blue"], bold=True, align=PP_ALIGN.CENTER, name="Consolas")
    # Right leg
    _rect(s, 3.4, 4.8, 0.5, 1.2, fill=RGBColor(0xAD, 0xD8, 0xE6))
    # Right foot
    _rect(s, 3.3, 6.0, 0.6, 0.4, fill=RGBColor(0x4A, 0x90, 0xC4),
          shape=MSO_SHAPE.OVAL)
    _add_textbox(s, 3.0, 6.35, 1.2, 0.3, "32.1\u00b0C",
                 size=9, color=C["blue"], bold=True, align=PP_ALIGN.CENTER, name="Consolas")

    # Right side: formula box
    _rounded_rect(s, 6.8, 1.8, 5.5, 0.8, fill=C["surface"], border_color=C["border"])
    _add_textbox(s, 7.0, 1.9, 5.1, 0.6, "CPTD  =  T_core  \u2212  T_peripheral",
                 size=22, color=C["teal"], bold=True, name="Consolas", align=PP_ALIGN.CENTER)

    _add_textbox(s, 6.8, 2.8, 5.5, 1.2,
                 "Core-to-peripheral temperature difference (CPTD) is a validated "
                 "marker of circulatory compromise. A rising CPTD indicates peripheral "
                 "vasoconstriction \u2014 the body\u2019s earliest autonomic response to "
                 "infection or shock.",
                 size=13, color=C["ink_light"])

    # 4 stat boxes
    stats = [
        ("5 hrs", "earlier than\nclinical signs", C["teal"], C["teal_light"]),
        ("83%", "sensitivity\n(Leante-Castellanos)", C["green"], C["green_light"]),
        ("0.3\u00b0C", "resolution\nFLIR Lepton 3.5", C["blue"], RGBColor(0xEB, 0xF2, 0xFA)),
        ("All tones", "works on\nevery skin tone", C["ink"], C["surface"]),
    ]
    sx = 6.8
    for val, lbl, nc, bg in stats:
        _small_stat_box(s, sx, 4.4, 1.25, 1.1, val, lbl, num_color=nc, bg=bg)
        sx += 1.35

    _add_textbox(s, 6.8, 5.8, 5.5, 0.8,
                 "Leante-Castellanos et al., 2015; Knobel-Dail et al., 2017; "
                 "Ussat et al., 2018 \u2014 see Appendix A5",
                 size=10, color=C["ink_muted"])

    # ---- Slide 7: Pipeline ---------------------------------------------------
    s = _blank(prs)
    _accent_bar(s)
    _label(s, 0.8, 0.3, "HOW IT WORKS")
    _title(s, 0.8, 0.7, "Four components. Existing infrastructure.")
    _subtitle(s, 0.8, 1.3, "No new network. No cloud. No custom EHR development.")

    boxes = [
        ("Incubator\nModule", "Custom-built device\nMounts in incubator lid\nFLIR Lepton + compute\n<$1,000 / bed",
         C["red"], C["red_light"]),
        ("On-Device AI", "Built into the module\nYOLOv8-Pose + CPTD\n24h local buffer\nNo separate hardware",
         C["amber"], RGBColor(0xFD, 0xF5, 0xE6)),
        ("InterSystems\nHealth Connect", "HL7v2 \u2194 FHIR R4\nMessage queue + audit\nRoutes to any EHR\nAlready at Northwell",
         C["teal"], C["teal_light"]),
        ("EHR Flowsheet", "CPTD as vitals\nBPA alerts\nWhere clinicians look\nEpic Wave 2: May 2026",
         C["blue"], RGBColor(0xEB, 0xF2, 0xFA)),
    ]
    bx = 0.6
    for i, (t, body, icon_clr, bg_clr) in enumerate(boxes):
        border_c = C["teal"] if i == 2 else C["border"]
        bw = 2 if i == 2 else 1
        _rounded_rect(s, bx, 2.2, 2.7, 3.6, fill=bg_clr, border_color=border_c, border_w=bw)
        # icon circle
        _rect(s, bx + 0.9, 2.4, 0.9, 0.7, fill=icon_clr, shape=MSO_SHAPE.OVAL)
        _add_textbox(s, bx + 0.15, 3.3, 2.4, 0.7, t,
                     size=16, color=C["ink"], bold=True, align=PP_ALIGN.CENTER)
        _add_textbox(s, bx + 0.15, 4.1, 2.4, 1.6, body,
                     size=12, color=C["ink_light"], align=PP_ALIGN.CENTER)
        # arrow
        if i < 3:
            _arrow_right(s, bx + 2.75, 3.8, 0.45, 0.25, fill=C["ink_muted"])
        bx += 3.15

    _add_textbox(s, 0.8, 6.2, 11.7, 0.5,
                 "NeoTherm deploys once. Health Connect routes everywhere.",
                 size=14, color=C["ink_muted"], align=PP_ALIGN.CENTER)

    # ---- Slide 8: InterSystems Integration -----------------------------------
    s = _blank(prs)
    _accent_bar(s, C["teal_deep"])
    _label(s, 0.8, 0.3, "INTERSYSTEMS INTEGRATION", C["teal_deep"])
    _title(s, 0.8, 0.7, "Health Connect is the backbone. FHIR R4 is the language.")

    # Left column: 3 stacked cards
    left_cards = [
        ("Phase 1 \u2014 HL7v2 via Capsule", "Edge box emits HL7v2 ORU. Capsule receives. "
         "Health Connect routes to EHR. Zero custom development.", C["teal_light"]),
        ("Phase 2\u20133 \u2014 FHIR R4 Direct", "Observations to Health Connect FHIR endpoint. "
         "IRIS repository. SMART on FHIR widget in Epic.", C["teal_light"]),
        ("EHR Migration Bridge", "Cohen Children\u2019s Epic go-live May 30, 2026. "
         "Health Connect routes to legacy before, Epic after. Deploy once.", C["surface"]),
    ]
    cy = 1.8
    for title_text, body_text, bg in left_cards:
        _rounded_rect(s, 0.8, cy, 5.8, 1.4, fill=bg, border_color=C["border"])
        _add_textbox(s, 1.0, cy + 0.15, 5.4, 0.4, title_text,
                     size=16, color=C["teal_deep"], bold=True)
        _add_textbox(s, 1.0, cy + 0.6, 5.4, 0.7, body_text,
                     size=13, color=C["ink_light"])
        cy += 1.55

    # Right column: dark FHIR JSON code block
    fhir_lines = [
        ('{', C["white"]),
        ('  "resourceType": "Observation",', RGBColor(0x98, 0xC3, 0x79)),
        ('  "status": "final",', RGBColor(0x98, 0xC3, 0x79)),
        ('  "code": {', C["white"]),
        ('    "coding": [{', C["white"]),
        ('      "system": "http://neotherm.io",', RGBColor(0xCE, 0x91, 0x78)),
        ('      "code": "cptd",', RGBColor(0xCE, 0x91, 0x78)),
        ('      "display": "Core-Peripheral', RGBColor(0xCE, 0x91, 0x78)),
        ('        Temperature Difference"', RGBColor(0xCE, 0x91, 0x78)),
        ('    }]', C["white"]),
        ('  },', C["white"]),
        ('  "valueQuantity": {', C["white"]),
        ('    "value": 2.7,', RGBColor(0xB5, 0xCE, 0xA8)),
        ('    "unit": "degC"', RGBColor(0xB5, 0xCE, 0xA8)),
        ('  },', C["white"]),
        ('  "interpretation": [{', C["white"]),
        ('    "coding": [{"code": "H"}]', RGBColor(0xFF, 0xD7, 0x00)),
        ('  }]', C["white"]),
        ('}', C["white"]),
    ]
    _code_block(s, 7.2, 1.8, 5.3, 5.2, fhir_lines)

    # ---- Slide 9: Demo Census ------------------------------------------------
    s = _blank(prs)
    _accent_bar(s)
    _label(s, 0.8, 0.3, "LIVE DEMO \u2014 NICU CENSUS")
    _title(s, 0.8, 0.7, "12 patients. NeoTherm flags visible at a glance.")
    try:
        _img(s, "census-view.png", 0.8, 1.8, width=11.7, height=4.8)
    except Exception:
        _rounded_rect(s, 0.8, 1.8, 11.7, 4.8, fill=C["surface"], border_color=C["border"])
        _add_textbox(s, 3, 3.5, 7, 1, "[census-view.png]", size=18, color=C["ink_muted"],
                     align=PP_ALIGN.CENTER)
    _add_textbox(s, 0.8, 6.8, 11.7, 0.4,
                 "Fully interactive Epic Hyperdrive replica. Built this weekend.",
                 size=13, color=C["ink_muted"], align=PP_ALIGN.CENTER)

    # ---- Slide 10: Demo BPA Alert --------------------------------------------
    s = _blank(prs)
    _accent_bar(s, C["red"])
    _label(s, 0.8, 0.3, "LIVE DEMO \u2014 BESTPRACTICE ALERT", C["red"])
    _title(s, 0.8, 0.7, "CPTD crosses threshold. Alert fires on chart open.")
    try:
        _img(s, "garcia-bpa-modal.png", 0.8, 1.8, width=11.7, height=4.8)
    except Exception:
        _rounded_rect(s, 0.8, 1.8, 11.7, 4.8, fill=C["surface"], border_color=C["border"])
    _add_textbox(s, 0.8, 6.8, 11.7, 0.4,
                 "Hospital-configured Epic BPA rule \u2014 no NeoTherm software generates the alert. No FDA clearance needed.",
                 size=13, color=C["ink_muted"], align=PP_ALIGN.CENTER)

    # ---- Slide 11: Demo NeoTherm Widget --------------------------------------
    s = _blank(prs)
    _accent_bar(s)
    _label(s, 0.8, 0.3, "LIVE DEMO \u2014 NEOTHERM WIDGET", C["teal"])
    _title(s, 0.8, 0.7, "Dedicated tab. Trend data. Zone breakdown.")
    try:
        _img(s, "garcia-neotherm.png", 0.8, 1.8, width=11.7, height=4.8)
    except Exception:
        _rounded_rect(s, 0.8, 1.8, 11.7, 4.8, fill=C["surface"], border_color=C["border"])
    _add_textbox(s, 0.8, 6.8, 11.7, 0.4,
                 "SMART on FHIR app queries InterSystems IRIS for Health. Deployed via Epic App Market.",
                 size=13, color=C["ink_muted"], align=PP_ALIGN.CENTER)

    # ---- Slide 12: Go-To-Market Timeline -------------------------------------
    s = _blank(prs)
    _accent_bar(s)
    _label(s, 0.8, 0.3, "GO-TO-MARKET")
    _title(s, 0.8, 0.7, "Built-out path to market.")

    milestones = [
        (C["blue"], "Cohen Children\u2019s Pilot", "Q4 2026",
         "1\u20132 beds. Display-only thermal + zone temps in Epic Flowsheet.\n"
         "IRB study with Dr. Barry Weinberger as Clinical PI. 200+ neonates target."),
        (C["teal"], "Silent Validation", "Q1\u2013Q2 2027",
         "Expand to 6\u201310 beds. CPTD runs silently \u2014 data recorded but alerts NOT shown to clinicians.\n"
         "Builds unbiased prospective dataset."),
        (C["phase3"], "Randomized Controlled Trial", "Q3 2027 \u2013 Q1 2028",
         "Full RCT: intervention arm gets alerts, control arm standard monitoring.\n"
         "Primary endpoint: time-to-antibiotic. Powers FDA De Novo."),
        (C["green"], "Commercial Launch", "2028+",
         "FDA clearance with RCT data. GPO contracts (Vizient, Premier).\n"
         "20 \u2192 100+ hospitals. Adult ICU expansion ($6\u20138B TAM)."),
    ]
    my = 1.8
    for bar_clr, m_title, timing, desc in milestones:
        _rect(s, 0.8, my, 0.12, 1.1, fill=bar_clr)
        _add_textbox(s, 1.1, my, 4.0, 0.4, m_title,
                     size=18, color=C["ink"], bold=True)
        _add_textbox(s, 5.5, my, 2.0, 0.4, timing,
                     size=14, color=bar_clr, bold=True)
        _add_textbox(s, 1.1, my + 0.4, 11.4, 0.7, desc,
                     size=13, color=C["ink_light"])
        my += 1.3

    # ---- Slide 13: Business Case ---------------------------------------------
    s = _blank(prs)
    _accent_bar(s)
    _label(s, 0.8, 0.3, "THE BUSINESS CASE")
    _title(s, 0.8, 0.7, "13.2x return. Pays for itself in 30 days.")

    # Left: 3 large stat boxes stacked
    _stat_box(s, 0.8, 1.9, 3.8, 1.2, "$82,500", "annual subscription per unit",
              num_color=C["teal"], bg=C["teal_light"], num_size=32)
    _stat_box(s, 0.8, 3.3, 3.8, 1.2, "$1.14M", "sepsis cost savings per unit",
              num_color=C["green"], bg=C["green_light"], num_size=32)
    _stat_box(s, 0.8, 4.7, 3.8, 1.2, "$65K", "CMS RPM reimbursement",
              num_color=C["blue"], bg=RGBColor(0xEB, 0xF2, 0xFA), num_size=32)

    # Right: ROI table
    _rounded_rect(s, 5.2, 1.9, 7.3, 3.2, fill=C["surface"], border_color=C["border"])
    rows = [
        ("Pays (subscription)", "$82,500 / yr", C["red"]),
        ("Gets (CMS RPM)", "+$31,987 / yr", C["green"]),
        ("Gets (sepsis savings)", "+$1,140,000 / yr", C["green"]),
    ]
    ry = 2.1
    for label_text, val, clr in rows:
        _add_textbox(s, 5.5, ry, 3.5, 0.35, label_text,
                     size=14, color=C["ink"])
        _add_textbox(s, 9.2, ry, 3.0, 0.35, val,
                     size=14, color=clr, bold=True, align=PP_ALIGN.RIGHT)
        ry += 0.55
    # separator
    _rect(s, 5.5, ry - 0.1, 6.5, 0.02, fill=C["border"])
    # Net benefit
    _rounded_rect(s, 5.4, ry + 0.05, 6.9, 0.55, fill=C["teal_light"], border_color=C["teal"])
    _add_textbox(s, 5.5, ry + 0.1, 3.5, 0.4, "Net benefit",
                 size=16, color=C["teal"], bold=True)
    _add_textbox(s, 9.2, ry + 0.1, 3.0, 0.4, "$1,089,487 / yr",
                 size=16, color=C["teal"], bold=True, align=PP_ALIGN.RIGHT)

    # Bottom metrics
    _small_stat_box(s, 5.5, 5.5, 2.1, 0.9, "13.2x", "ROI",
                    num_color=C["teal"], bg=C["teal_light"])
    _small_stat_box(s, 7.8, 5.5, 2.1, 0.9, "30 days", "payback period",
                    num_color=C["green"], bg=C["green_light"])
    _small_stat_box(s, 10.1, 5.5, 2.1, 0.9, "82%", "gross margin",
                    num_color=C["blue"], bg=RGBColor(0xEB, 0xF2, 0xFA))

    # ---- Slide 14: Market ----------------------------------------------------
    s = _blank(prs)
    _accent_bar(s)
    _label(s, 0.8, 0.3, "MARKET")
    _title(s, 0.8, 0.7, "$2.1 billion problem. Zero competitors. We are first.")

    # Left: red card
    _rounded_rect(s, 0.8, 1.8, 5.5, 5.0, fill=C["red_light"], border_color=C["red"])
    _add_textbox(s, 1.1, 1.95, 5.0, 0.5, "$2.14B",
                 size=36, color=C["red"], bold=True)
    _add_textbox(s, 1.1, 2.5, 5.0, 0.4, "Annual US Cost Burden",
                 size=16, color=C["red"])

    # TAM / SAM / SOM funnel
    funnel = [
        (4.8, "TAM   $1.6B", "All US NICU beds (16,000+)", C["red"]),
        (3.8, "SAM   $400M", "Level III\u2013IV NICUs (4,800 beds)", C["amber"]),
        (2.8, "SOM   $33M", "First 400 beds by 2030", C["green"]),
    ]
    fy = 3.3
    for fw, label_text, desc, clr in funnel:
        fx = 1.0 + (4.8 - fw) / 2
        _rounded_rect(s, fx, fy, fw, 0.7, fill=clr)
        _add_textbox(s, fx + 0.1, fy + 0.05, fw - 0.2, 0.3, label_text,
                     size=14, color=C["white"], bold=True, align=PP_ALIGN.CENTER)
        _add_textbox(s, fx + 0.1, fy + 0.35, fw - 0.2, 0.3, desc,
                     size=10, color=C["white"], align=PP_ALIGN.CENTER)
        fy += 0.85

    # Right: dark card
    _rounded_rect(s, 6.8, 1.8, 5.7, 5.0, fill=C["ink"])
    _add_textbox(s, 7.2, 2.2, 4.8, 1.0, "0",
                 size=72, color=C["teal"], bold=True, align=PP_ALIGN.CENTER)
    _add_textbox(s, 7.2, 3.4, 4.8, 0.6, "commercial thermal camera\nproducts for NICU",
                 size=16, color=C["white"], align=PP_ALIGN.CENTER)
    _add_textbox(s, 7.2, 4.2, 4.8, 1.0,
                 "\u201CNo commercial system exists for continuous non-contact "
                 "thermal monitoring in the NICU.\u201D\n\u2014 Frontiers in Pediatrics, 2023",
                 size=12, color=RGBColor(0xA0, 0xAE, 0xBF), italic=True,
                 align=PP_ALIGN.CENTER)
    _add_textbox(s, 7.2, 5.5, 4.8, 0.8,
                 "Adult ICU expansion\n$6\u20138B TAM",
                 size=18, color=C["teal"], bold=True, align=PP_ALIGN.CENTER)

    # ---- Slide 15: Northwell Partnership -------------------------------------
    s = _blank(prs)
    _accent_bar(s)
    _label(s, 0.8, 0.3, "NORTHWELL PARTNERSHIP")
    _title(s, 0.8, 0.7, "Champions. Infrastructure. Ready.")

    # Left: 2 stakeholder cards
    _card(s, 0.8, 1.8, 5.5, 2.2, top_border_color=C["teal"])
    _add_textbox(s, 1.0, 1.95, 5.1, 0.3, "Clinical Champion",
                 size=11, color=C["teal"], bold=True)
    _add_textbox(s, 1.0, 2.25, 5.1, 0.5, "Barry Weinberger, MD",
                 size=24, color=C["ink"], name="Georgia")
    _add_textbox(s, 1.0, 2.75, 5.1, 0.4,
                 "Assoc. Director, Neonatology \u2014 Cohen Children\u2019s",
                 size=13, color=C["ink_light"])
    _add_textbox(s, 1.0, 3.15, 5.1, 0.4, "NICU access  \u00b7  Clinical PI  \u00b7  IRB",
                 size=13, color=C["teal"], bold=True)

    _card(s, 0.8, 4.2, 5.5, 2.2, top_border_color=C["blue"])
    _add_textbox(s, 1.0, 4.35, 5.1, 0.3, "Innovation Champion",
                 size=11, color=C["blue"], bold=True)
    _add_textbox(s, 1.0, 4.65, 5.1, 0.5, "Theodoros Zanos, PhD",
                 size=24, color=C["ink"], name="Georgia")
    _add_textbox(s, 1.0, 5.15, 5.1, 0.4,
                 "AVP, Exploration & Innovation \u2014 Health AI",
                 size=13, color=C["ink_light"])
    _add_textbox(s, 1.0, 5.55, 5.1, 0.4,
                 "16-person team  \u00b7  IT approvals  \u00b7  Health Connect",
                 size=13, color=C["blue"], bold=True)

    # Right: Infrastructure table
    _rounded_rect(s, 6.8, 1.8, 5.7, 3.5, fill=C["surface"], border_color=C["border"])
    _add_textbox(s, 7.0, 1.9, 5.3, 0.35, "Infrastructure Already In Place",
                 size=15, color=C["ink"], bold=True)
    infra_rows = [
        ("InterSystems Health Connect", "Message broker (live)"),
        ("Capsule MDI", "Device integration (live)"),
        ("Epic Hyperdrive", "Go-live May 30, 2026"),
        ("FHIR R4 Endpoint", "Via Health Connect"),
        ("Network / Wi-Fi", "Existing NICU infra"),
    ]
    iy = 2.35
    for sys, role in infra_rows:
        _add_textbox(s, 7.1, iy, 2.8, 0.3, sys,
                     size=12, color=C["ink"], bold=True)
        _add_textbox(s, 9.9, iy, 2.4, 0.3, role,
                     size=12, color=C["ink_light"])
        iy += 0.42

    # Teal note
    _rounded_rect(s, 6.8, 5.6, 5.7, 0.8, fill=C["teal_light"], border_color=C["teal"])
    _add_textbox(s, 7.0, 5.65, 5.3, 0.7,
                 "Fully reversible: NeoTherm adds to existing infrastructure. "
                 "Unplug and the NICU returns to normal.",
                 size=12, color=C["teal_deep"])

    # ---- Slide 16: Regulatory Roadmap ----------------------------------------
    s = _blank(prs)
    _accent_bar(s)
    _label(s, 0.8, 0.3, "REGULATORY ROADMAP")
    _title(s, 0.8, 0.7, "Phase 1 needs no FDA. We start now.")

    phases = [
        ("Phase 1", C["blue"], "Display & Record",
         "NO FDA SUBMISSION",
         ["\u2022 Thermal image + zone temps only",
          "\u2022 Displayed as \u201Cadjunctive\u201D data",
          "\u2022 No diagnostic claim",
          "\u2022 Clinician interprets all data",
          "\u2022 IRB oversight only"]),
        ("Phase 2", C["teal"], "CPTD Trending",
         "ENFORCEMENT DISCRETION",
         ["\u2022 Calculated CPTD metric",
          "\u2022 Threshold-based alerting",
          "\u2022 Clinical Decision Support",
          "\u2022 FDA CDS exclusion criteria",
          "\u2022 Transparent algorithm"]),
        ("Phase 3", C["phase3"], "AI Sepsis Warning",
         "510(k) / DE NOVO",
         ["\u2022 ML-based risk prediction",
          "\u2022 Autonomous alert generation",
          "\u2022 Requires clinical trial data",
          "\u2022 RCT powers submission",
          "\u2022 Predicate: InSightIQ (De Novo)"]),
    ]
    px = 0.8
    for phase_name, clr, phase_title, badge, bullets in phases:
        # header
        _rounded_rect(s, px, 1.8, 3.8, 0.6, fill=clr)
        _add_textbox(s, px + 0.1, 1.85, 3.6, 0.5,
                     f"{phase_name}: {phase_title}",
                     size=15, color=C["white"], bold=True, align=PP_ALIGN.CENTER)
        # badge
        _rounded_rect(s, px + 0.3, 2.55, 3.2, 0.45, fill=C["surface"], border_color=clr, border_w=2)
        _add_textbox(s, px + 0.4, 2.57, 3.0, 0.4, badge,
                     size=12, color=clr, bold=True, align=PP_ALIGN.CENTER)
        # bullets
        by = 3.15
        for b in bullets:
            _add_textbox(s, px + 0.2, by, 3.4, 0.35, b,
                         size=12, color=C["ink_light"])
            by += 0.38
        px += 4.05

    # Progress bar
    _rect(s, 0.8, 6.4, 3.8, 0.15, fill=C["blue"])
    _rect(s, 4.65, 6.4, 3.8, 0.15, fill=C["teal"])
    _rect(s, 8.5, 6.4, 3.8, 0.15, fill=C["phase3"])
    _add_textbox(s, 0.8, 6.6, 3.8, 0.3, "NOW", size=10, color=C["blue"],
                 bold=True, align=PP_ALIGN.CENTER)
    _add_textbox(s, 4.65, 6.6, 3.8, 0.3, "2027", size=10, color=C["teal"],
                 bold=True, align=PP_ALIGN.CENTER)
    _add_textbox(s, 8.5, 6.6, 3.8, 0.3, "2028+", size=10, color=C["phase3"],
                 bold=True, align=PP_ALIGN.CENTER)

    # ---- Slide 17: InterSystems Challenge ------------------------------------
    s = _blank(prs)
    _accent_bar(s, C["teal_deep"])
    _label(s, 0.8, 0.3, "THE INTERSYSTEMS CHALLENGE", C["teal_deep"])
    _title(s, 0.8, 0.7, "Why Global-Therma wins.")

    # 2x2 grid
    grid = [
        (0.8, 1.8, "Real Interoperability",
         "Not a wrapper. Health Connect is the production message broker at Northwell \u2014 "
         "17,000+ clinicians, 23 hospitals."),
        (5.2, 1.8, "HL7v2 + FHIR R4",
         "Day-one HL7v2 via Capsule for immediate deployment. FHIR R4 for modern "
         "EHR integration. Both through Health Connect."),
        (0.8, 3.7, "EHR Migration Ready",
         "Cohen Children\u2019s transitions from Allscripts Sunrise to Epic May 2026. "
         "Health Connect routes messages to either. Deploy once."),
        (5.2, 3.7, "IRIS Repository",
         "CPTD history stored in InterSystems IRIS for Health. Powers trend "
         "analysis, ML training, and SMART on FHIR queries."),
    ]
    for gx, gy, gt, gb in grid:
        _rounded_rect(s, gx, gy, 4.0, 1.6, fill=C["teal_light"], border_color=C["border"])
        _add_textbox(s, gx + 0.2, gy + 0.15, 3.6, 0.35, gt,
                     size=16, color=C["teal_deep"], bold=True)
        _add_textbox(s, gx + 0.2, gy + 0.55, 3.6, 0.9, gb,
                     size=12, color=C["ink_light"])

    # Right sidebar: dark checklist
    _rounded_rect(s, 9.8, 1.8, 3.0, 3.5, fill=C["ink"])
    checks = ["\u2713 Health Connect", "\u2713 FHIR R4 Facade", "\u2713 IRIS for Health",
              "\u2713 SMART on FHIR", "\u2713 Message queuing", "\u2713 Multi-EHR routing"]
    cy = 1.95
    for ck in checks:
        _add_textbox(s, 10.0, cy, 2.6, 0.35, ck,
                     size=14, color=C["teal"], bold=True)
        cy += 0.48

    # Below sidebar: teal banner
    _rounded_rect(s, 9.8, 5.5, 3.0, 1.0, fill=C["teal"])
    _add_textbox(s, 9.9, 5.55, 2.8, 0.9,
                 "Not a prototype.\nProduction architecture at the largest US health system.",
                 size=12, color=C["white"], bold=True, align=PP_ALIGN.CENTER)

    # ---- Slide 18: The Ask ---------------------------------------------------
    s = _blank(prs)
    _accent_bar(s)
    _add_textbox(s, 1.0, 1.0, 11.333, 0.5,
                 "680,000 newborns die from circulatory failure every year\u2026",
                 size=16, color=C["ink_muted"], align=PP_ALIGN.CENTER)
    _add_textbox(s, 1.0, 2.0, 11.333, 1.2,
                 "We just need to plug it in.",
                 size=54, color=C["ink"], name="Georgia", align=PP_ALIGN.CENTER)

    # 4 stats
    ask_stats = [
        ("5 hrs", "earlier detection"),
        ("13.2x", "ROI"),
        ("$0 infra", "new infrastructure"),
        ("0", "competitors"),
    ]
    asx = 1.8
    for val, lbl in ask_stats:
        _add_textbox(s, asx, 3.6, 2.3, 0.6, val,
                     size=36, color=C["teal"], bold=True, align=PP_ALIGN.CENTER)
        _add_textbox(s, asx, 4.2, 2.3, 0.4, lbl,
                     size=13, color=C["ink_muted"], align=PP_ALIGN.CENTER)
        asx += 2.5

    _add_textbox(s, 1.0, 5.0, 11.333, 0.5,
                 "One camera. InterSystems Health Connect. Every NICU in America.",
                 size=18, color=C["ink_light"], align=PP_ALIGN.CENTER)
    _add_textbox(s, 1.0, 5.8, 11.333, 0.5,
                 "Global-Therma. Early detection saves lives.",
                 size=22, color=C["ink"], bold=True, align=PP_ALIGN.CENTER)

    # ---- Slide 19: Blue Closer -----------------------------------------------
    s = _blank(prs)
    _solid_bg(s, C["hack_blue"])
    _add_textbox(s, 1, 0.8, 11.333, 0.4, "TEAM 218",
                 size=14, color=RGBColor(0x66, 0x96, 0xD6), align=PP_ALIGN.CENTER)
    _add_textbox(s, 1, 1.4, 11.333, 0.8, "Global-Therma",
                 size=44, color=C["white"], bold=True, align=PP_ALIGN.CENTER)

    team = [
        ("Diego Pedregal", "CEO / Product"),
        ("Marcos Rodriguez", "CTO / Engineering"),
        ("Rida Baquai", "Clinical / Neonatology"),
        ("Michael Meehan", "Business Dev"),
        ("Taylir Speer", "Regulatory / Strategy"),
        ("Jedd Horowitz", "InterSystems / Infra"),
    ]
    # 2 columns, 3 rows
    for i, (name, role) in enumerate(team):
        col = i % 2
        row = i // 2
        x = 2.5 + col * 4.5
        y = 2.8 + row * 0.65
        _add_textbox(s, x, y, 2.2, 0.35, name,
                     size=16, color=C["white"], bold=True)
        _add_textbox(s, x + 2.3, y, 2.0, 0.35, role,
                     size=14, color=RGBColor(0xA0, 0xBE, 0xE0))

    _add_textbox(s, 1, 5.8, 11.333, 0.5, "\u2713  InterSystems Prize",
                 size=18, color=C["white"], bold=True, align=PP_ALIGN.CENTER)

    # Save
    out = os.path.join(BASE, "deck.pptx")
    prs.save(out)
    print(f"Saved {out}  ({len(prs.slides)} slides)")
    return out


# ===========================================================================
# APPENDIX  (9 slides)
# ===========================================================================

def build_appendix():
    prs = _new_prs()

    # ---- A1: Full Architecture Diagram ---------------------------------------
    s = _blank(prs)
    _accent_bar(s)
    _label(s, 0.8, 0.3, "APPENDIX A1 \u2014 FULL ARCHITECTURE DIAGRAM")
    _title(s, 0.8, 0.7, "End-to-end data flow: device to EHR.", size=30)

    cols = [
        ("Edge / Incubator", [
            ("FLIR Lepton 3.5", "160\u00d7120 thermal sensor\n9 Hz capture", "P1\u20133"),
            ("Coral Edge TPU", "YOLOv8-Pose inference\nOn-device processing", "P1\u20133"),
            ("Local Buffer", "24h circular buffer\nSQLite + raw frames", "P1\u20133"),
        ]),
        ("Transport", [
            ("HL7v2 ORU", "Observation Result\nUnsolicited message", "P1"),
            ("FHIR R4 REST", "Observation resource\nHTTPS POST", "P2\u20133"),
            ("Capsule MDI", "Medical Device\nIntegration gateway", "P1"),
        ]),
        ("Integration", [
            ("Health Connect", "Message broker\nRouting + transform", "P1\u20133"),
            ("IRIS for Health", "FHIR repository\nClinical data store", "P2\u20133"),
            ("Audit / Queue", "Message persistence\nCompliance logging", "P1\u20133"),
        ]),
        ("EHR / Clinical", [
            ("Epic Flowsheet", "CPTD as vital sign\nZone temperatures", "P1\u20133"),
            ("BPA Alert", "Best Practice Advisory\nThreshold-triggered", "P2\u20133"),
            ("SMART on FHIR", "NeoTherm widget\nTrend + heatmap", "P3"),
        ]),
    ]
    cx = 0.4
    for col_title, items in cols:
        _add_textbox(s, cx, 1.5, 3.0, 0.4, col_title,
                     size=14, color=C["teal_deep"], bold=True, align=PP_ALIGN.CENTER)
        iy = 2.0
        for name, desc, phase in items:
            _rounded_rect(s, cx, iy, 2.9, 1.3, fill=C["surface"], border_color=C["border"])
            _add_textbox(s, cx + 0.1, iy + 0.08, 2.7, 0.3, name,
                         size=13, color=C["ink"], bold=True)
            _add_textbox(s, cx + 0.1, iy + 0.4, 2.7, 0.55, desc,
                         size=10, color=C["ink_light"])
            # Phase badge
            phase_clr = C["blue"] if "1" in phase else C["teal"]
            _rounded_rect(s, cx + 1.9, iy + 0.95, 0.8, 0.25, fill=phase_clr)
            _add_textbox(s, cx + 1.9, iy + 0.95, 0.8, 0.25, phase,
                         size=8, color=C["white"], bold=True, align=PP_ALIGN.CENTER)
            iy += 1.45
        # arrows to next column
        if cx < 9:
            _arrow_right(s, cx + 3.0, 3.2, 0.3, 0.2, fill=C["ink_muted"])
        cx += 3.25

    # ---- A2: HL7v2 Message Format --------------------------------------------
    s = _blank(prs)
    _accent_bar(s)
    _label(s, 0.8, 0.3, "APPENDIX A2 \u2014 HL7v2 MESSAGE FORMAT")
    _title(s, 0.8, 0.7, "ORU^R01 observation message for CPTD transmission.", size=30)

    hl7_lines = [
        ("MSH|^~\\&|NEOTHERM|NICU_BED_3|", RGBColor(0xFF, 0xD7, 0x00)),
        ("  HEALTHCONNECT|NORTHWELL|", RGBColor(0xFF, 0xD7, 0x00)),
        ("  20260815143022||ORU^R01|", RGBColor(0x98, 0xC3, 0x79)),
        ("  MSG00001|P|2.5.1", RGBColor(0x98, 0xC3, 0x79)),
        ("PID|1||MRN12345^^^MRN||", RGBColor(0xCE, 0x91, 0x78)),
        ("  GARCIA^ELENA||20260101|F", RGBColor(0xCE, 0x91, 0x78)),
        ("OBR|1||ORD001|CPTD^Core-Periph", RGBColor(0xB5, 0xCE, 0xA8)),
        ("  Temp Diff^NEOTHERM", RGBColor(0xB5, 0xCE, 0xA8)),
        ("OBX|1|NM|CPTD^Core-Peripheral", RGBColor(0x9C, 0xDC, 0xFE)),
        ("  Temperature Difference^NEOTHERM", RGBColor(0x9C, 0xDC, 0xFE)),
        ("  ||2.7|degC|0.0-2.0|H|||F", RGBColor(0x9C, 0xDC, 0xFE)),
        ("OBX|2|NM|TCORE^Core Temp", RGBColor(0xD4, 0xA5, 0xFF)),
        ("  ^NEOTHERM||36.8|degC||||F", RGBColor(0xD4, 0xA5, 0xFF)),
        ("OBX|3|NM|TPERIPH^Peripheral", RGBColor(0xD4, 0xA5, 0xFF)),
        ("  Temp^NEOTHERM||34.1|degC||||F", RGBColor(0xD4, 0xA5, 0xFF)),
    ]
    _code_block(s, 0.8, 1.6, 7.5, 5.0, hl7_lines)

    # Info boxes
    info = [
        ("Transport", "TCP/IP via\nCapsule MDI", C["blue"]),
        ("Routing", "Health Connect\nmessage broker", C["teal"]),
        ("Frequency", "Every 60 seconds\n(configurable)", C["green"]),
    ]
    ibx = 8.8
    for it, ib, ic in info:
        _rounded_rect(s, ibx, 1.6, 3.8, 1.3, fill=C["surface"], border_color=ic, border_w=2)
        _add_textbox(s, ibx + 0.2, 1.7, 3.4, 0.3, it,
                     size=14, color=ic, bold=True)
        _add_textbox(s, ibx + 0.2, 2.1, 3.4, 0.6, ib,
                     size=12, color=C["ink_light"])
        ibx += 0
        break  # stack vertically
    sby = 3.1
    for it, ib, ic in info[1:]:
        _rounded_rect(s, 8.8, sby, 3.8, 1.3, fill=C["surface"], border_color=ic, border_w=2)
        _add_textbox(s, 9.0, sby + 0.1, 3.4, 0.3, it,
                     size=14, color=ic, bold=True)
        _add_textbox(s, 9.0, sby + 0.5, 3.4, 0.6, ib,
                     size=12, color=C["ink_light"])
        sby += 1.5

    # ---- A3: FHIR R4 Observation ---------------------------------------------
    s = _blank(prs)
    _accent_bar(s)
    _label(s, 0.8, 0.3, "APPENDIX A3 \u2014 FHIR R4 OBSERVATION")
    _title(s, 0.8, 0.7, "Observation resource for CPTD.", size=30)

    fhir_full = [
        ('{', C["white"]),
        ('  "resourceType": "Observation",', RGBColor(0x98, 0xC3, 0x79)),
        ('  "id": "cptd-20260815-143022",', RGBColor(0x98, 0xC3, 0x79)),
        ('  "status": "final",', RGBColor(0x98, 0xC3, 0x79)),
        ('  "category": [{"coding": [{"system":', RGBColor(0xCE, 0x91, 0x78)),
        ('    "http://terminology.hl7.org/', RGBColor(0xCE, 0x91, 0x78)),
        ('    CodeSystem/observation-category",', RGBColor(0xCE, 0x91, 0x78)),
        ('    "code": "vital-signs"}]}],', RGBColor(0xCE, 0x91, 0x78)),
        ('  "code": {"coding": [{"system":', C["white"]),
        ('    "http://neotherm.io/fhir/codes",', RGBColor(0x9C, 0xDC, 0xFE)),
        ('    "code": "cptd",', RGBColor(0x9C, 0xDC, 0xFE)),
        ('    "display": "Core-Peripheral', RGBColor(0x9C, 0xDC, 0xFE)),
        ('      Temperature Difference"}]},', RGBColor(0x9C, 0xDC, 0xFE)),
        ('  "subject": {"reference":', C["white"]),
        ('    "Patient/garcia-elena"},', RGBColor(0xB5, 0xCE, 0xA8)),
        ('  "effectiveDateTime":', C["white"]),
        ('    "2026-08-15T14:30:22Z",', RGBColor(0xB5, 0xCE, 0xA8)),
        ('  "valueQuantity": {', C["white"]),
        ('    "value": 2.7,', RGBColor(0xFF, 0xD7, 0x00)),
        ('    "unit": "degC",', RGBColor(0xFF, 0xD7, 0x00)),
        ('    "system": "http://unitsofmeasure', RGBColor(0xFF, 0xD7, 0x00)),
        ('      .org", "code": "Cel"},', RGBColor(0xFF, 0xD7, 0x00)),
        ('  "interpretation": [{"coding":', C["white"]),
        ('    [{"code": "H",', RGBColor(0xD4, 0xA5, 0xFF)),
        ('    "display": "High"}]}]', RGBColor(0xD4, 0xA5, 0xFF)),
        ('}', C["white"]),
    ]
    _code_block(s, 0.8, 1.6, 11.7, 5.6, fhir_full)

    # ---- A4: Regulatory Decision Table ---------------------------------------
    s = _blank(prs)
    _accent_bar(s)
    _label(s, 0.8, 0.3, "APPENDIX A4 \u2014 REGULATORY DECISION TABLE")
    _title(s, 0.8, 0.7, "FDA classification analysis for each phase.", size=30)

    # Table
    table_data = [
        ("Question", "Answer", "Source"),
        ("Is NeoTherm a medical device?", "Phase 1: No (display only)\nPhase 2-3: Yes", "21 CFR 820"),
        ("Does Phase 1 make a diagnosis?", "No. Displays temperatures only.", "FDA CDS Guidance 2022"),
        ("Is CPTD a new biomarker?", "No. Validated in literature since 1991.", "Leante-Castellanos 2015"),
        ("Does Phase 2 qualify for CDS exclusion?", "Yes, if transparent + clinician review.", "21st Century Cures \u00a73060"),
        ("What class is Phase 3?", "Class II (510(k) or De Novo)", "FDA Product Code QMT"),
        ("Is there a predicate device?", "InSightIQ (DEN200063) \u2014 De Novo", "FDA De Novo Database"),
        ("Is the camera FDA-cleared?", "FLIR Lepton is non-medical; system-level clearance needed.", "FDA 510(k) Database"),
        ("What about off-label thermal use?", "78% of NICU devices used off-label (AAP 2015).", "AAP Technical Report"),
        ("Who is the regulatory contact?", "Taylir Speer \u2014 regulatory strategy lead.", "Internal"),
    ]
    rows_count = len(table_data)
    cols_count = 3
    tbl_shape = s.shapes.add_table(rows_count, cols_count,
                                   Inches(0.8), Inches(1.6),
                                   Inches(11.7), Inches(5.4))
    tbl = tbl_shape.table
    col_widths = [Inches(3.5), Inches(5.0), Inches(3.2)]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w

    for r in range(rows_count):
        for c_idx in range(cols_count):
            cell = tbl.cell(r, c_idx)
            cell.text = table_data[r][c_idx]
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(11)
                    run.font.name = "Calibri"
                    if r == 0:
                        run.font.bold = True
                        run.font.color.rgb = C["white"]
                    else:
                        run.font.color.rgb = C["ink"] if c_idx < 2 else C["ink_muted"]
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["teal"]
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["surface"]

    # ---- A5: Clinical Evidence -----------------------------------------------
    s = _blank(prs)
    _accent_bar(s)
    _label(s, 0.8, 0.3, "APPENDIX A5 \u2014 CLINICAL EVIDENCE")
    _title(s, 0.8, 0.7, "Published evidence supporting CPTD as a sepsis biomarker.", size=30)

    evidence = [
        ("Claim", "Evidence", "Source"),
        ("CPTD predicts sepsis 5h early",
         "Core-peripheral temp gradient widened 5.2h before clinical diagnosis.",
         "Leante-Castellanos 2015"),
        ("83% sensitivity for LOS",
         "CPTD >2\u00b0C had 83% sensitivity, 75% specificity for late-onset sepsis.",
         "Leante-Castellanos 2015"),
        ("Infrared matches contact probes",
         "Non-contact IR agreed within 0.2\u00b0C of skin contact sensors.",
         "Knobel-Dail 2017"),
        ("CPTD correlates with NEC",
         "Peripheral cooling preceded NEC diagnosis by 6\u201312 hours.",
         "Ussat 2018"),
        ("Works across all skin tones",
         "Thermal IR emissivity 0.97\u20130.98 regardless of melanin content.",
         "FLIR Technical Note 2020"),
        ("Continuous > intermittent",
         "Continuous monitoring detected 3x more hypothermic episodes.",
         "Knobel-Dail 2017"),
        ("Cost of missed sepsis",
         "Each sepsis episode adds $170,000 in NICU costs.",
         "Swanson 2018; Rolnitsky 2023"),
        ("CPTD validated in premature + term",
         "Studies include 24\u201340 weeks GA across multiple centers.",
         "Multiple (see references)"),
    ]
    rows_count = len(evidence)
    cols_count = 3
    tbl_shape = s.shapes.add_table(rows_count, cols_count,
                                   Inches(0.8), Inches(1.6),
                                   Inches(11.7), Inches(5.4))
    tbl = tbl_shape.table
    col_widths = [Inches(3.2), Inches(5.5), Inches(3.0)]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w

    for r in range(rows_count):
        for c_idx in range(cols_count):
            cell = tbl.cell(r, c_idx)
            cell.text = evidence[r][c_idx]
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(11)
                    run.font.name = "Calibri"
                    if r == 0:
                        run.font.bold = True
                        run.font.color.rgb = C["white"]
                    else:
                        run.font.color.rgb = C["ink"]
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["teal"]
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["surface"]

    # ---- A6: Contract Stack --------------------------------------------------
    s = _blank(prs)
    _accent_bar(s)
    _label(s, 0.8, 0.3, "APPENDIX A6 \u2014 CONTRACT STACK")
    _title(s, 0.8, 0.7, "Nine agreements required before pilot deployment.", size=30)

    contracts = [
        ("BAA", "Business Associate\nAgreement", "HIPAA-required for\nany PHI access", C["red"]),
        ("IRB Protocol", "Institutional Review\nBoard Approval", "Display-only study\nExempt or Expedited", C["blue"]),
        ("IT Security", "Northwell InfoSec\nReview", "Network scan, pentest\nData flow diagram", C["ink_muted"]),
        ("MSA", "Master Services\nAgreement", "Commercial terms\nLiability, IP, term", C["teal"]),
        ("DUA", "Data Use\nAgreement", "De-identified data\nfor ML training", C["amber"]),
        ("Clinical Eng.", "Biomed Engineering\nApproval", "Electrical safety\nIEC 60601 testing", C["green"]),
        ("E&O Insurance", "Errors & Omissions\nProfessional Liability", "$2M per occurrence\n$5M aggregate", C["phase3"]),
        ("Cyber Insurance", "Cybersecurity\nLiability Coverage", "Required by\nNorthwell IT policy", C["red"]),
        ("Product Liability", "General Liability\n+ Product Coverage", "Non-diagnostic claim\nreduces exposure", C["blue"]),
    ]
    gx = 0.6
    gy = 1.6
    for i, (title_text, desc, detail, clr) in enumerate(contracts):
        if i > 0 and i % 3 == 0:
            gx = 0.6
            gy += 1.8
        _card(s, gx, gy, 3.9, 1.6, top_border_color=clr, top_border_h=0.06)
        _add_textbox(s, gx + 0.15, gy + 0.15, 3.6, 0.3, title_text,
                     size=16, color=C["ink"], bold=True)
        _add_textbox(s, gx + 0.15, gy + 0.5, 3.6, 0.5, desc,
                     size=11, color=C["ink_light"])
        _add_textbox(s, gx + 0.15, gy + 1.05, 3.6, 0.5, detail,
                     size=10, color=C["ink_muted"])
        gx += 4.1

    # Adjunctive tool note
    _rounded_rect(s, 0.6, 6.6, 12.1, 0.55, fill=C["teal_light"], border_color=C["teal"])
    _add_textbox(s, 0.8, 6.65, 11.7, 0.4,
                 "Adjunctive Tool Classification: Phase 1 is display-only, reducing insurance and liability requirements significantly.",
                 size=12, color=C["teal_deep"], bold=True, align=PP_ALIGN.CENTER)

    # ---- A7: Northwell Infrastructure ----------------------------------------
    s = _blank(prs)
    _accent_bar(s)
    _label(s, 0.8, 0.3, "APPENDIX A7 \u2014 NORTHWELL INFRASTRUCTURE")
    _title(s, 0.8, 0.7, "Existing systems at Cohen Children\u2019s Medical Center.", size=30)

    infra_data = [
        ("System", "Role", "Source"),
        ("InterSystems Health Connect", "HL7v2 / FHIR message broker, routing engine, audit log.",
         "Northwell IT (confirmed by Zanos team)"),
        ("Capsule MDI (Qualcomm)", "Medical device integration gateway for bedside monitors.",
         "Northwell Biomed Engineering"),
        ("Allscripts Sunrise (current)", "Current EHR at Cohen Children\u2019s. Transitions to Epic May 2026.",
         "Northwell CIO Office"),
        ("Epic Hyperdrive (incoming)", "New EHR. Flowsheet, BPA, SMART on FHIR support.",
         "Epic Wave 2 deployment plan"),
        ("InterSystems IRIS for Health", "FHIR R4 data repository. Powers analytics + SMART apps.",
         "Northwell Health AI team"),
    ]
    rows_count = len(infra_data)
    cols_count = 3
    tbl_shape = s.shapes.add_table(rows_count, cols_count,
                                   Inches(0.8), Inches(1.6),
                                   Inches(11.7), Inches(3.5))
    tbl = tbl_shape.table
    col_widths = [Inches(3.5), Inches(5.0), Inches(3.2)]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w

    for r in range(rows_count):
        for c_idx in range(cols_count):
            cell = tbl.cell(r, c_idx)
            cell.text = infra_data[r][c_idx]
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12)
                    run.font.name = "Calibri"
                    if r == 0:
                        run.font.bold = True
                        run.font.color.rgb = C["white"]
                    else:
                        run.font.color.rgb = C["ink"]
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["teal"]
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["surface"]

    # EHR transition note
    _rounded_rect(s, 0.8, 5.4, 11.7, 1.0, fill=C["teal_light"], border_color=C["teal"])
    _add_textbox(s, 1.0, 5.5, 11.3, 0.8,
                 "EHR Transition: Cohen Children\u2019s moves from Allscripts Sunrise to Epic Hyperdrive "
                 "on May 30, 2026. Health Connect abstracts this transition \u2014 NeoTherm deploys once "
                 "and messages route to whichever EHR is active. This is a core InterSystems value proposition.",
                 size=12, color=C["teal_deep"])

    # ---- A8: Hardware Cost ---------------------------------------------------
    s = _blank(prs)
    _accent_bar(s)
    _label(s, 0.8, 0.3, "APPENDIX A8 \u2014 HARDWARE COST BREAKDOWN")
    _title(s, 0.8, 0.7, "Bill of materials: prototype to scale.", size=30)

    bom = [
        ("Component", "Prototype Cost", "At Scale (1000+)"),
        ("FLIR Lepton 3.5 (160\u00d7120)", "$250", "$120"),
        ("Coral Edge TPU (USB)", "$60", "$35"),
        ("Raspberry Pi CM4 (4GB)", "$65", "$40"),
        ("Custom PCB + enclosure", "$200", "$45"),
        ("Power supply + cabling", "$30", "$15"),
        ("Assembly + QC", "$150", "$20"),
        ("TOTAL", "$755", "$275"),
    ]
    rows_count = len(bom)
    cols_count = 3
    tbl_shape = s.shapes.add_table(rows_count, cols_count,
                                   Inches(0.8), Inches(1.6),
                                   Inches(7.0), Inches(3.5))
    tbl = tbl_shape.table
    col_widths = [Inches(3.5), Inches(1.75), Inches(1.75)]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w

    for r in range(rows_count):
        for c_idx in range(cols_count):
            cell = tbl.cell(r, c_idx)
            cell.text = bom[r][c_idx]
            for p in cell.text_frame.paragraphs:
                if c_idx > 0:
                    p.alignment = PP_ALIGN.RIGHT
                for run in p.runs:
                    run.font.size = Pt(12)
                    run.font.name = "Calibri"
                    if r == 0:
                        run.font.bold = True
                        run.font.color.rgb = C["white"]
                    elif r == rows_count - 1:
                        run.font.bold = True
                        run.font.color.rgb = C["teal"]
                    else:
                        run.font.color.rgb = C["ink"]
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["teal"]
            elif r == rows_count - 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["teal_light"]
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["surface"]

    # Scale curve table
    scale_data = [
        ("Volume", "Unit Cost", "Margin"),
        ("1 (prototype)", "$755", "N/A"),
        ("10 (pilot)", "$500", "78%"),
        ("100", "$350", "86%"),
        ("1,000+", "$275", "90%+"),
    ]
    rows_count2 = len(scale_data)
    tbl2_shape = s.shapes.add_table(rows_count2, 3,
                                    Inches(8.5), Inches(1.6),
                                    Inches(4.5), Inches(2.5))
    tbl2 = tbl2_shape.table
    for i, w in enumerate([Inches(1.8), Inches(1.35), Inches(1.35)]):
        tbl2.columns[i].width = w

    for r in range(rows_count2):
        for c_idx in range(3):
            cell = tbl2.cell(r, c_idx)
            cell.text = scale_data[r][c_idx]
            for p in cell.text_frame.paragraphs:
                if c_idx > 0:
                    p.alignment = PP_ALIGN.RIGHT
                for run in p.runs:
                    run.font.size = Pt(11)
                    run.font.name = "Calibri"
                    if r == 0:
                        run.font.bold = True
                        run.font.color.rgb = C["white"]
                    else:
                        run.font.color.rgb = C["ink"]
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["blue"]
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["surface"]

    # 3 highlight boxes
    _stat_box(s, 0.8, 5.5, 3.2, 1.1, "<$1K", "prototype unit",
              num_color=C["red"], bg=C["red_light"])
    _stat_box(s, 4.3, 5.5, 3.2, 1.1, "$275", "at 1,000+ scale",
              num_color=C["teal"], bg=C["teal_light"])
    _stat_box(s, 7.8, 5.5, 3.2, 1.1, "90%+", "gross margin at scale",
              num_color=C["green"], bg=C["green_light"])

    # ---- A9: LOS Cost Savings ------------------------------------------------
    s = _blank(prs)
    _accent_bar(s)
    _label(s, 0.8, 0.3, "APPENDIX A9 \u2014 LENGTH-OF-STAY COST SAVINGS")
    _title(s, 0.8, 0.7, "Conservative model: earlier detection reduces LOS and cost.", size=30)

    # Left: Two source cards
    _card(s, 0.8, 1.8, 5.5, 1.8, top_border_color=C["blue"])
    _add_textbox(s, 1.0, 1.95, 5.1, 0.3, "Source: Swanson et al. 2018",
                 size=13, color=C["blue"], bold=True)
    _add_textbox(s, 1.0, 2.3, 5.1, 1.1,
                 "\u2022 Neonatal sepsis episodes cost avg $170,927 more than matched controls\n"
                 "\u2022 Average additional LOS: 22.4 days\n"
                 "\u2022 Based on 41,666 NICU admissions across 348 US hospitals",
                 size=12, color=C["ink_light"])

    _card(s, 0.8, 3.8, 5.5, 1.8, top_border_color=C["teal"])
    _add_textbox(s, 1.0, 3.95, 5.1, 0.3, "Source: Rolnitsky et al. 2023",
                 size=13, color=C["teal"], bold=True)
    _add_textbox(s, 1.0, 4.3, 5.1, 1.1,
                 "\u2022 Earlier empiric antibiotics reduced sepsis mortality by 23%\n"
                 "\u2022 Each hour of delay increased mortality odds by 4%\n"
                 "\u2022 Earlier treatment correlated with 18-25% shorter sepsis LOS",
                 size=12, color=C["ink_light"])

    # Right: Calculation code block
    calc_lines = [
        ("# NeoTherm LOS Savings Model", RGBColor(0x6A, 0x99, 0x55)),
        ("", C["white"]),
        ("sepsis_cost    = $170,927", RGBColor(0xB5, 0xCE, 0xA8)),
        ("sepsis_los     = 22.4 days", RGBColor(0xB5, 0xCE, 0xA8)),
        ("cost_per_day   = $7,630", RGBColor(0xB5, 0xCE, 0xA8)),
        ("", C["white"]),
        ("los_reduction  = 20%  (conservative)", RGBColor(0xFF, 0xD7, 0x00)),
        ("days_saved     = 4.48 days", RGBColor(0xFF, 0xD7, 0x00)),
        ("savings_per_ep = $34,185", RGBColor(0xFF, 0xD7, 0x00)),
        ("", C["white"]),
        ("nicu_beds      = 50", RGBColor(0x9C, 0xDC, 0xFE)),
        ("sepsis_rate    = 25%", RGBColor(0x9C, 0xDC, 0xFE)),
        ("admissions/yr  = 800", RGBColor(0x9C, 0xDC, 0xFE)),
        ("episodes/yr    = 200", RGBColor(0x9C, 0xDC, 0xFE)),
        ("", C["white"]),
        ("TOTAL SAVINGS  = $6,837,000/yr", RGBColor(0x98, 0xC3, 0x79)),
        ("PER BED        = $136,740/yr", RGBColor(0x98, 0xC3, 0x79)),
    ]
    _code_block(s, 6.8, 1.8, 5.7, 4.8, calc_lines)

    # 3 result boxes
    _stat_box(s, 0.8, 5.9, 3.5, 1.1, "20%", "LOS reduction (conservative)",
              num_color=C["teal"], bg=C["teal_light"])
    _stat_box(s, 4.6, 5.9, 3.5, 1.1, "$6,591", "saved per patient",
              num_color=C["green"], bg=C["green_light"])
    _stat_box(s, 8.4, 5.9, 3.5, 1.1, "$213M", "US-wide annual savings potential",
              num_color=C["blue"], bg=RGBColor(0xEB, 0xF2, 0xFA))

    # Save
    out = os.path.join(BASE, "appendix.pptx")
    prs.save(out)
    print(f"Saved {out}  ({len(prs.slides)} slides)")
    return out


# ===========================================================================
if __name__ == "__main__":
    build_deck()
    build_appendix()
    print("\nDone. Both files generated.")
